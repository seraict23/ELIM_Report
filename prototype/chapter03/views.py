from django.shortcuts import render

# Create your views here.

from .models import StatusTable
from common.models import Common

from pathlib import Path
import win32com.client as win32
import win32api

from threading import Thread


# work Thread
import sys
import os
sys.path.append(os.path.dirname(os.path.abspath(os.path.dirname(__file__))))
from engine.chapter3 import Worker



# 파워포인트
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Cm
import os
def MartToPPTX(xPoint, yPoint):
    PATH = os.getcwd()

    prs = Presentation(PATH+'/engine/1.pptx')

    slide = prs.slides[0]

    left=Cm(xPoint)
    top=Cm(yPoint)
    pic = slide.shapes.add_picture(PATH+"/engine/pointer.jpg", left, top)

    prs.save(PATH+'/engine/2.pptx')



def Chapter3(req):
    statustable = StatusTable.objects.all().order_by('id')

    context = {
        'statustable':statustable
    }

    return render(req, 'chapter03/chapter03.html', context)


def ExcelData(req):
    currentBuilding = Common.objects.get(pk=1)
    fileName = req.POST['excelFile']

    # 엑셀 처리
    BASE_DIR = Path(__file__).resolve().parent.parent
    PATH = str(BASE_DIR)+"/src/"

    class ExcelApp:
        def __init__(self):
            self.excel = win32.Dispatch("Excel.Application")
            self.excel.Visible = True

            wb = self.excel.WorkBooks.Open(PATH+fileName)
            ws = wb.Sheets("결함리스트")

            # 원하는 범위만큼 셀 읽기

            count = 0
            for row in range(10, 3000):
                valueList = []
                for col in range(1, 21):
                    valueList.append(str(ws.Cells(row, col).Value))
                if valueList.count("None")>10:
                    print(count, " 종료")
                    break
                else : 
                    print(valueList)
                    photo_number = 'No.'+valueList[14].split('.')[0] if valueList[14]!="None" else ""
                    number = valueList[2].split('.')[0] if valueList[2]!="None" else ""
                    StatusTable.objects.create(
                        building = currentBuilding,
                        location = valueList[1] if valueList[1]!="None" else "",
                        number = number,
                        part = valueList[5],
                        material = valueList[6],
                        shapeType = valueList[7],
                        width = valueList[8],
                        LengthArea = valueList[9],
                        each = valueList[10],
                        progress = valueList[11],
                        note = valueList[12],
                        cause = valueList[13],
                        photoNumber =photo_number,
                        photoName=valueList[18] if valueList[18]!="None" else "",
                    )
                count += 1
            win32api.MessageBox(0, "완료되었습니다. 새로고침을 누르세요", "확인창", 0)

    def createExcel():
        import pythoncom
        pythoncom.CoInitialize()
        excel = ExcelApp()
        pythoncom.CoUninitialize()

    thread = Thread(target = createExcel)
    thread.start()

    statustable = StatusTable.objects.all().order_by('id')

    context = {
        'statustable':statustable
    }

    return render(req, 'chapter03/chapter03.html', context)


def ClearData(req):
    all = StatusTable.objects.all()
    all.delete()

    statustable = StatusTable.objects.all().order_by('id')

    context = {
        'statustable':statustable
    }
    return render(req, 'chapter03/chapter03.html', context) 


def ViewData(req):
    statustable = StatusTable.objects.all().order_by('id')

    queryset = [e for e in statustable]
    list2d = []
    for i in queryset:
        listA = []
        listA.append(i.location)
        listA.append(i.number)
        listA.append(i.part)
        listA.append(i.material)
        listA.append(i.shapeType)
        listA.append(i.width)
        listA.append(i.LengthArea)
        listA.append(i.each)
        listA.append(i.progress)
        listA.append(i.note)
        listA.append(i.cause)
        listA.append(i.photoNumber)
        listA.append(i.photoName)
        list2d.append(listA)
    
    # print(list2d)

    sub = Worker(list2d)
    sub.start()

    context = {
        'statustable':statustable
    }
    return render(req, 'chapter03/chapter03.html', context) 


def Drawing(req):
    if req.method == "POST":
        xPoint = float(req.POST['x'])
        xPoint = (xPoint-50)/960*25.4

        yPoint = float(req.POST['y'])
        yPoint = (yPoint-71)/720*19.05

        print(xPoint)
        print(yPoint)

        MartToPPTX(xPoint,yPoint)

        return render(req, 'chapter03/drawing.html')
    else :

        return render(req, 'chapter03/drawing.html')
    
