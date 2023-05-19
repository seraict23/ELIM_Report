import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Cm

import os

PATH = os.getcwd()

prs = Presentation(PATH+'/1.pptx')

slide = prs.slides[0]

left=Cm(14.9)
top=Cm(10.5)
pic = slide.shapes.add_picture("./pointer.jpg", left, top)

prs.save(PATH+'/2.pptx')